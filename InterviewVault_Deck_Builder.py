"""
InterviewVault — Advanced PowerPoint Deck Builder
==================================================
Generates a polished, dark-themed .pptx covering:
  - Title, Problem Statements, Vision
  - System Architecture (high-level + backend + frontend + AI layer)
  - Methodologies (adaptive learning path, diagnostics, Socratic follow-ups,
    company intelligence, skill profiling)
  - Problems We Solved
  - Tech Stack
  - Data Model
  - Body Language & Communication Analysis
  - Roadmap & Differentiators
  - Closing slide

Prerequisites
-------------
    pip install python-pptx

Run
---
    python InterviewVault_Deck_Builder.py

Output
------
    InterviewVault_Master_Deck.pptx (16:9, dark theme, indigo/purple accents)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Brand palette  (matches the project's CSS variables)
# ---------------------------------------------------------------------------
BG_DEEP        = RGBColor(0x05, 0x05, 0x0A)   # background void
BG_PANEL       = RGBColor(0x0E, 0x0E, 0x1A)   # card panel
BG_PANEL_2     = RGBColor(0x14, 0x14, 0x24)   # secondary panel
PRIMARY        = RGBColor(0x63, 0x66, 0xF1)   # indigo
ACCENT         = RGBColor(0xA8, 0x55, 0xF7)   # purple
SUCCESS        = RGBColor(0x10, 0xB9, 0x81)   # emerald (Green list)
WARNING        = RGBColor(0xF5, 0x9E, 0x0B)   # amber  (Yellow list)
DANGER         = RGBColor(0xEF, 0x44, 0x44)
TEXT           = RGBColor(0xF1, 0xF5, 0xF9)
MUTED          = RGBColor(0x94, 0xA3, 0xB8)
DIM            = RGBColor(0x64, 0x74, 0x8B)
BORDER         = RGBColor(0x1F, 0x1F, 0x33)

HEAD_FONT = "Calibri"          # change to "Space Grotesk" if installed locally
BODY_FONT = "Calibri"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
BLANK   = prs.slide_layouts[6]


def add_bg(slide, color=BG_DEEP):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w if line_w else 0.75)
    shp.shadow.inherit = False
    return shp


def add_round(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=TEXT, font=BODY_FONT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_multi(slide, x, y, w, h, runs, *, size=14, color=TEXT, font=BODY_FONT,
              align=PP_ALIGN.LEFT, line_spacing=1.25):
    """runs = list of (text, dict-of-overrides) or list of str."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, item in enumerate(runs):
        if isinstance(item, str):
            text, opts = item, {}
        else:
            text, opts = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        r = p.add_run()
        r.text = text
        r.font.name = opts.get("font", font)
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", color)
    return tb


def add_header(slide, eyebrow, title, *, y=0.55):
    add_text(slide, Inches(0.7), Inches(y), Inches(12), Inches(0.35),
             eyebrow.upper(), size=11, bold=True, color=ACCENT, font=HEAD_FONT)
    add_text(slide, Inches(0.7), Inches(y + 0.3), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=TEXT, font=HEAD_FONT,
             line_spacing=1.05)
    # underline accent block (small, premium — not a horizontal rule)
    add_rect(slide, Inches(0.7), Inches(y + 1.05), Inches(0.45), Inches(0.06),
             fill=PRIMARY)


def add_footer(slide, idx, total, section="InterviewVault"):
    add_text(slide, Inches(0.7), Inches(7.05), Inches(6), Inches(0.3),
             section, size=10, color=DIM, font=HEAD_FONT)
    add_text(slide, Inches(11.7), Inches(7.05), Inches(1.0), Inches(0.3),
             f"{idx:02d} / {total:02d}", size=10, color=DIM, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, *, fill=BG_PANEL, accent=PRIMARY, accent_w=0.07):
    """Glass card with a thin left accent bar."""
    add_rect(slide, x, y, w, h, fill=fill, line=BORDER, line_w=0.75)
    add_rect(slide, x, y, Inches(accent_w), h, fill=accent)


def chip(slide, x, y, label, *, color=PRIMARY, text_color=TEXT, width=None):
    text = f"  {label}  "
    w = width if width is not None else Inches(0.06 * len(label) + 0.5)
    h = Inches(0.32)
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.adjustments[0] = 0.5
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    tf = bg.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = HEAD_FONT
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = text_color
    return w


def section_band(slide, color=PRIMARY):
    """Tiny top brand band."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.06), fill=color)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
TOTAL = 16


def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG_DEEP)

    # subtle large brand orb (decoration)
    orb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5),
                             Inches(6), Inches(6))
    orb.fill.solid(); orb.fill.fore_color.rgb = PRIMARY
    orb.line.fill.background()
    orb.shadow.inherit = False
    # fake transparency by overlaying a near-bg rect with low opacity is not
    # supported cleanly in python-pptx; we use a darker overlay instead
    add_rect(s, Inches(9.5), Inches(-1.5), Inches(6), Inches(6),
             fill=BG_DEEP).fill.fore_color.rgb = BG_DEEP

    orb2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.4), Inches(-0.6),
                              Inches(4.2), Inches(4.2))
    orb2.fill.solid(); orb2.fill.fore_color.rgb = ACCENT
    orb2.line.fill.background(); orb2.shadow.inherit = False
    add_rect(s, Inches(10.4), Inches(-0.6), Inches(4.2), Inches(4.2),
             fill=BG_DEEP)

    # Brand mark
    add_rect(s, Inches(0.7), Inches(0.7), Inches(0.45), Inches(0.06),
             fill=PRIMARY)
    add_text(s, Inches(0.7), Inches(0.85), Inches(8), Inches(0.35),
             "INTERVIEWVAULT", size=11, bold=True, color=ACCENT,
             font=HEAD_FONT)

    # Hero
    add_text(s, Inches(0.7), Inches(2.0), Inches(11.8), Inches(1.3),
             "The Udemy for Interview Prep —",
             size=44, bold=True, color=TEXT, font=HEAD_FONT,
             line_spacing=1.05)
    add_text(s, Inches(0.7), Inches(2.85), Inches(11.8), Inches(1.3),
             "powered by GenAI & Agentic AI.",
             size=44, bold=True, color=PRIMARY, font=HEAD_FONT,
             line_spacing=1.05)

    # Subtitle
    add_text(s, Inches(0.7), Inches(4.1), Inches(11), Inches(0.6),
             "Adaptive, gamified, role-aware interview preparation that "
             "emulates a real end-to-end interview.",
             size=16, color=MUTED, font=BODY_FONT, line_spacing=1.35)

    # Chips
    x0 = Inches(0.7); y0 = Inches(5.0)
    labels = ["Adaptive Diagnostics", "Socratic Follow-ups",
              "Company Intelligence", "Body-Language Aware",
              "Per-User Skill Vector"]
    for lab in labels:
        w = chip(s, x0, y0, lab, color=BG_PANEL_2, text_color=TEXT)
        x0 += w + Inches(0.12)

    # Footer
    add_text(s, Inches(0.7), Inches(6.85), Inches(8), Inches(0.3),
             "Advanced technical overview  ·  System architecture · "
             "Methodologies · Tech stack",
             size=11, color=DIM)
    add_text(s, Inches(0.7), Inches(7.15), Inches(8), Inches(0.3),
             "Prepared by Nitin   ·   2026",
             size=10, color=DIM, italic=True)


# ---------------------------------------------------------------------------
# Slide 2 — Problem Statements
# ---------------------------------------------------------------------------
def slide_problem_statements():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s, PRIMARY)
    add_header(s, "01 · Problem Statements",
               "Interview prep is broken — and it's costing real careers.")

    problems = [
        ("Generic content overload",
         "LeetCode + YouTube + blogs ≠ a personalized path. Learners "
         "drown in content with no signal on what they actually need.",
         DANGER),
        ("No feedback loop",
         "Practice without diagnosis. Users repeat the same mistakes "
         "for weeks and don't know which topics are dragging them down.",
         WARNING),
        ("Mock interviews are theatre",
         "Most tools fire random questions. They ignore the candidate's "
         "weaknesses, the role, and the target company's interview style.",
         PRIMARY),
        ("Non-technical signals ignored",
         "Body language, eye contact, voice pace, filler words and "
         "fluency all matter — but no platform measures them.",
         ACCENT),
        ("Company-specific prep is manual",
         "Glassdoor → blog → LinkedIn → Reddit. Hours of scavenging, "
         "with no synthesis into an actionable study plan.",
         SUCCESS),
        ("Time-pressure scenarios unsupported",
         "Interview in 24 hours? In a week? Existing tools don't "
         "re-rank or compress the syllabus to the time available.",
         RGBColor(0x06, 0xB6, 0xD4)),
    ]
    cols = 3
    cw = Inches(3.95)
    ch = Inches(2.1)
    x0, y0 = Inches(0.7), Inches(2.4)
    gx, gy = Inches(0.15), Inches(0.18)
    for i, (h, body, c) in enumerate(problems):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        card(s, x, y, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.25), y + Inches(0.2), cw - Inches(0.35),
                 Inches(0.4), h, size=14, bold=True, color=TEXT,
                 font=HEAD_FONT)
        add_text(s, x + Inches(0.25), y + Inches(0.6), cw - Inches(0.35),
                 Inches(1.4), body, size=11, color=MUTED, line_spacing=1.35)

    add_footer(s, 2, TOTAL, "Problem Statements")


# ---------------------------------------------------------------------------
# Slide 3 — Vision & Core Philosophy
# ---------------------------------------------------------------------------
def slide_vision():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "02 · Vision", "What InterviewVault actually is.")

    # Left: big statement
    add_text(s, Inches(0.7), Inches(2.5), Inches(7), Inches(2.5),
             "An adaptive, role-aware coach that gets smarter every time you engage with it.",
             size=26, bold=True, color=TEXT, font=HEAD_FONT, line_spacing=1.15)
    add_text(s, Inches(0.7), Inches(4.4), Inches(7), Inches(2.2),
             "Every user gets a personalized syllabus, diagnostic-driven "
             "weaknesses, Socratic follow-ups, company-specific prep, and "
             "an interview that mirrors the real thing — technical, "
             "behavioral, communication, and body language.",
             size=14, color=MUTED, line_spacing=1.45)

    # Right: 4 pillars stacked
    pillars = [
        ("Personalized", "Path tailored to role + diagnostic weak spots.", PRIMARY),
        ("Adaptive",     "Difficulty + topics rebalance after every activity.", ACCENT),
        ("Holistic",     "Tech + behavioral + body + voice scored together.", SUCCESS),
        ("Compounding",  "Skill vector persists; the system learns the user.", WARNING),
    ]
    x = Inches(8.1); y = Inches(2.3); w = Inches(4.5); h = Inches(1.05)
    for i, (t, sub, c) in enumerate(pillars):
        card(s, x, y + i * (h + Inches(0.12)), w, h, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.3), y + i * (h + Inches(0.12)) + Inches(0.18),
                 w - Inches(0.4), Inches(0.4), t,
                 size=15, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.3), y + i * (h + Inches(0.12)) + Inches(0.55),
                 w - Inches(0.4), Inches(0.45), sub,
                 size=11, color=MUTED, line_spacing=1.3)

    add_footer(s, 3, TOTAL, "Vision")


# ---------------------------------------------------------------------------
# Slide 4 — High-Level System Architecture
# ---------------------------------------------------------------------------
def slide_system_arch():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "03 · System Architecture",
               "A four-tier, AI-orchestrated platform.")

    # 4 tiers stacked horizontally as columns
    tiers = [
        ("CLIENT",
         ["React 18 + Vite", "shadcn/ui + Radix", "Framer Motion",
          "@dnd-kit", "Recharts", "TanStack Query"],
         PRIMARY),
        ("API / SERVICES",
         ["FastAPI (Python)", "JWT Auth", "Role-aware routers",
          "Assessment / Interview / Diagnostic", "Onboarding & Path",
          "Skill Profile service"],
         ACCENT),
        ("AI ORCHESTRATION",
         ["Gemini 2.5 Flash (JSON mode)", "Perplexity Sonar (live web)",
          "Whisper (ASR)", "face-api.js (gaze + landmarks)",
          "Prompt registry", "5-pass JSON parser"],
         SUCCESS),
        ("DATA",
         ["User + UserSkillProfile", "InterviewSession (transcript+report)",
          "CompanyInsights cache", "TopicArticle cache",
          "DiagnosticResult", "Behavioral stats (extensible)"],
         WARNING),
    ]
    cw = Inches(3.05); ch = Inches(4.0)
    x0, y0 = Inches(0.7), Inches(2.3); gap = Inches(0.12)
    for i, (h, items, c) in enumerate(tiers):
        x = x0 + i * (cw + gap)
        card(s, x, y0, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.25), y0 + Inches(0.2),
                 cw - Inches(0.4), Inches(0.4),
                 h, size=12, bold=True, color=c, font=HEAD_FONT)
        for j, it in enumerate(items):
            yy = y0 + Inches(0.65) + j * Inches(0.45)
            add_rect(s, x + Inches(0.28), yy + Inches(0.13),
                     Inches(0.09), Inches(0.09), fill=c)
            add_text(s, x + Inches(0.45), yy, cw - Inches(0.6), Inches(0.4),
                     it, size=11, color=TEXT, line_spacing=1.2)

    # Bottom arrow strip — data flow
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             "Data flow:  User action  →  FastAPI route  →  AI orchestration "
             "(Gemini/Perplexity)  →  Skill-profile update  →  UI re-render",
             size=11, italic=True, color=MUTED, font=BODY_FONT)

    add_footer(s, 4, TOTAL, "System Architecture")


# ---------------------------------------------------------------------------
# Slide 5 — Backend Architecture
# ---------------------------------------------------------------------------
def slide_backend():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "04 · Backend Architecture",
               "Service-oriented Python core, AI as a first-class layer.")

    # Two-column: left = router map, right = services
    add_text(s, Inches(0.7), Inches(2.4), Inches(5.8), Inches(0.4),
             "ROUTERS  (FastAPI)", size=12, bold=True, color=ACCENT,
             font=HEAD_FONT)
    routes = [
        ("/api/users",            "Auth, profile, target_role, onboarding"),
        ("/api/onboarding",       "Role pick, resume OCR, path init"),
        ("/api/learning-path",    "Green/Yellow lists, drag-drop persistence"),
        ("/api/diagnostic",       "Progressive difficulty Q gen + eval"),
        ("/api/articles",         "Topic articles (cached)"),
        ("/api/assessments",      "Quizzes + Socratic follow-up endpoint"),
        ("/api/interview",        "Session orchestration + transcript store"),
        ("/api/company-insights", "Perplexity → Gemini synthesis (cached)"),
        ("/api/skill-profile",    "Read/update per-topic 0–100 vector"),
    ]
    y = Inches(2.85)
    for path, desc in routes:
        card(s, Inches(0.7), y, Inches(5.8), Inches(0.42),
             fill=BG_PANEL, accent=PRIMARY, accent_w=0.05)
        add_text(s, Inches(0.85), y + Inches(0.06), Inches(2.2), Inches(0.32),
                 path, size=11, bold=True, color=TEXT,
                 font="Consolas")
        add_text(s, Inches(3.1), y + Inches(0.06), Inches(3.3), Inches(0.32),
                 desc, size=10, color=MUTED)
        y += Inches(0.46)

    # Right: services
    add_text(s, Inches(6.9), Inches(2.4), Inches(6.0), Inches(0.4),
             "SERVICES", size=12, bold=True, color=ACCENT, font=HEAD_FONT)
    services = [
        ("ai_service",
         "Gemini wrapper, 5-pass JSON parser, prompt templates."),
        ("perplexity_service",
         "Sonar calls for live company / role retrieval."),
        ("skill_profile_service",
         "Updates topic mastery 0–100 + history after every activity."),
        ("interview_engine",
         "Picks next question from weak-topic queue + diagnostic results."),
        ("behavioral_engine",
         "Post-hoc voice + body analysis from stored stats."),
        ("path_seeder",
         "Hard-coded Green paths per role + Gemini-generated Yellow set."),
    ]
    y = Inches(2.85)
    for name, desc in services:
        card(s, Inches(6.9), y, Inches(6.0), Inches(0.65),
             fill=BG_PANEL_2, accent=SUCCESS, accent_w=0.05)
        add_text(s, Inches(7.05), y + Inches(0.07),
                 Inches(2.5), Inches(0.3), name,
                 size=11, bold=True, color=TEXT, font="Consolas")
        add_text(s, Inches(7.05), y + Inches(0.35),
                 Inches(5.7), Inches(0.3),
                 desc, size=10, color=MUTED)
        y += Inches(0.72)

    add_footer(s, 5, TOTAL, "Backend Architecture")


# ---------------------------------------------------------------------------
# Slide 6 — Frontend Architecture
# ---------------------------------------------------------------------------
def slide_frontend():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "05 · Frontend Architecture",
               "A modern React stack engineered for *feel*.")

    # 2 rows × 4 cards
    items = [
        ("shadcn/ui + Radix", "Accessible primitives; consistent dark "
         "theme via CSS vars (--dk-*).", PRIMARY),
        ("Framer Motion", "Page transitions, layoutId shared-element, "
         "spring micro-interactions.", ACCENT),
        ("@dnd-kit", "Smooth drag-and-drop for the Green/Yellow topic "
         "configurator.", SUCCESS),
        ("TanStack Query", "Server state, cache invalidation, optimistic "
         "updates after activities.", WARNING),

        ("Recharts + SVG", "Skill radar, history charts, conic-gradient "
         "skill rings.", PRIMARY),
        ("cmdk + react-hot-toast", "Power-user command palette + clean "
         "toast notifications.", ACCENT),
        ("react-markdown", "Renders Gemini-generated articles with GFM "
         "support.", SUCCESS),
        ("Vaul + resizable-panels", "Drawer & split-pane layouts (article + "
         "notes side-by-side).", WARNING),
    ]
    cols = 4
    cw = Inches(3.0); ch = Inches(1.85)
    x0, y0 = Inches(0.7), Inches(2.4); gx, gy = Inches(0.12), Inches(0.15)
    for i, (t, d, c) in enumerate(items):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(s, x, y, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.22), y + Inches(0.18),
                 cw - Inches(0.35), Inches(0.4), t,
                 size=13, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.22), y + Inches(0.55),
                 cw - Inches(0.35), Inches(1.2), d,
                 size=10.5, color=MUTED, line_spacing=1.35)

    add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.4),
             "Design language: dark-first (#05050A), glassmorphic cards, "
             "indigo→purple gradients, Space Grotesk headings.",
             size=11, italic=True, color=MUTED)

    add_footer(s, 6, TOTAL, "Frontend Architecture")


# ---------------------------------------------------------------------------
# Slide 7 — AI / Agentic Orchestration
# ---------------------------------------------------------------------------
def slide_ai_layer():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "06 · AI Orchestration",
               "Right model for the right job — Gemini + Perplexity.")

    # Two columns: Gemini | Perplexity
    def model_card(x, y, w, h, name, role, jobs, c):
        card(s, x, y, w, h, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.3), y + Inches(0.25), w - Inches(0.4),
                 Inches(0.45), name, size=20, bold=True, color=TEXT,
                 font=HEAD_FONT)
        add_text(s, x + Inches(0.3), y + Inches(0.7), w - Inches(0.4),
                 Inches(0.35), role, size=12, italic=True, color=c,
                 font=HEAD_FONT)
        yy = y + Inches(1.2)
        for j in jobs:
            add_rect(s, x + Inches(0.32), yy + Inches(0.12),
                     Inches(0.08), Inches(0.08), fill=c)
            add_text(s, x + Inches(0.5), yy, w - Inches(0.7),
                     Inches(0.4), j, size=11.5, color=TEXT,
                     line_spacing=1.3)
            yy += Inches(0.42)

    model_card(Inches(0.7), Inches(2.3), Inches(6.0), Inches(4.4),
               "Gemini 2.5 Flash",
               "Structured generation · reasoning · scoring",
               ["Diagnostic question gen (easy → adv ladder)",
                "Topic article generation (cached)",
                "Multi-dimensional answer scoring",
                "Socratic follow-ups (vague / wrong / good / extend)",
                "Voice transcript → fluency, vocab, grammar",
                "Synthesizes Perplexity results into a study plan"],
               PRIMARY)

    model_card(Inches(7.0), Inches(2.3), Inches(5.6), Inches(4.4),
               "Perplexity Sonar",
               "Live web retrieval with citations",
               ["Company × role interview pattern lookup",
                "Glassdoor / LeetCode / eng-blog citations",
                "Always fresh — no stale training-cutoff data",
                "Cached into CompanyInsights table on first hit",
                "Fallback path when Gemini hits rate limits"],
               ACCENT)

    add_footer(s, 7, TOTAL, "AI Orchestration")


# ---------------------------------------------------------------------------
# Slide 8 — Methodology: Adaptive Learning Path
# ---------------------------------------------------------------------------
def slide_methodology_path():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "07 · Methodology",
               "Adaptive Learning Path — Green / Yellow + Diagnostic.")

    # Diagram: 3 stages
    stages = [
        ("01 · Role + Resume",
         "Pick target role. Optional resume → OCR → Gemini "
         "extracts skills, suggests role-fit cards.", PRIMARY),
        ("02 · Path Configurator",
         "Hard-coded Green path for the role + Gemini-generated "
         "Yellow extensions. Drag-drop to commit.", ACCENT),
        ("03 · Diagnostic Ladder",
         "Per topic: easy → intermediate → advanced. Stop on first "
         "fail. Classification feeds back into Green/Yellow.", SUCCESS),
    ]
    cw = Inches(4.0); ch = Inches(2.6)
    x0, y0 = Inches(0.7), Inches(2.3); gap = Inches(0.2)
    for i, (h, d, c) in enumerate(stages):
        x = x0 + i * (cw + gap)
        card(s, x, y0, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.3), y0 + Inches(0.25),
                 cw - Inches(0.5), Inches(0.45), h,
                 size=16, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.3), y0 + Inches(0.85),
                 cw - Inches(0.5), Inches(1.6), d,
                 size=12, color=MUTED, line_spacing=1.4)

    # Output legend
    legend_y = Inches(5.2)
    add_text(s, Inches(0.7), legend_y, Inches(12), Inches(0.4),
             "Diagnostic outcome → list assignment",
             size=12, bold=True, color=ACCENT, font=HEAD_FONT)
    legend = [("≤ Easy",  "Weak",         "Green (must study)", SUCCESS),
              ("Intermediate", "Mid",      "Green (reinforce)",  PRIMARY),
              ("Advanced", "Expert",      "Yellow (optional)",   WARNING)]
    yy = legend_y + Inches(0.45)
    xx = Inches(0.7)
    for level, tag, dest, c in legend:
        card(s, xx, yy, Inches(4.0), Inches(0.85), fill=BG_PANEL_2,
             accent=c, accent_w=0.05)
        add_text(s, xx + Inches(0.2), yy + Inches(0.1),
                 Inches(3.6), Inches(0.3),
                 f"Score: {level}", size=11, bold=True, color=c)
        add_text(s, xx + Inches(0.2), yy + Inches(0.4),
                 Inches(3.6), Inches(0.3),
                 f"{tag}  →  {dest}", size=10.5, color=TEXT)
        xx += Inches(4.15)

    add_footer(s, 8, TOTAL, "Methodology · Path")


# ---------------------------------------------------------------------------
# Slide 9 — Methodology: Socratic Follow-up + Skill Profile
# ---------------------------------------------------------------------------
def slide_methodology_socratic():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "08 · Methodology",
               "Socratic follow-ups + per-user skill vector.")

    # Left: Socratic loop
    add_text(s, Inches(0.7), Inches(2.4), Inches(5.8), Inches(0.4),
             "SOCRATIC FOLLOW-UP LOOP", size=12, bold=True, color=ACCENT,
             font=HEAD_FONT)

    loop = [
        ("Base question", "Pre-generated from topic + difficulty.", PRIMARY),
        ("Student answer", "Captured + length-checked (≥20 chars).", PRIMARY),
        ("Strategy classifier",
         "Vague / Wrong / Good / Extend → routes to one of 4 prompts.",
         ACCENT),
        ("Gemini follow-up",
         "ONE contextual probe based on the actual answer.", ACCENT),
        ("Score + update",
         "Both base + follow-up answers are scored and feed back "
         "into UserSkillProfile.", SUCCESS),
    ]
    y = Inches(2.85)
    for i, (h, d, c) in enumerate(loop):
        card(s, Inches(0.7), y, Inches(5.8), Inches(0.7),
             fill=BG_PANEL, accent=c, accent_w=0.05)
        add_text(s, Inches(0.95), y + Inches(0.08), Inches(0.4),
                 Inches(0.55), f"{i+1}", size=20, bold=True, color=c,
                 font=HEAD_FONT)
        add_text(s, Inches(1.35), y + Inches(0.08),
                 Inches(2.0), Inches(0.32), h,
                 size=12, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, Inches(1.35), y + Inches(0.38),
                 Inches(5.0), Inches(0.3),
                 d, size=10, color=MUTED, line_spacing=1.3)
        y += Inches(0.78)

    # Right: skill profile
    add_text(s, Inches(6.9), Inches(2.4), Inches(6.0), Inches(0.4),
             "USER SKILL PROFILE  (persistent vector)",
             size=12, bold=True, color=ACCENT, font=HEAD_FONT)

    add_text(s, Inches(6.9), Inches(2.85), Inches(6.0), Inches(0.45),
             "{user_id, topic, score 0–100, confidence, last_updated, history[]}",
             size=11, color=TEXT, font="Consolas")

    # Sample radar-like text + bars
    topics = [("Statistics", 78, SUCCESS),
              ("ML Fundamentals", 62, PRIMARY),
              ("NLP", 41, WARNING),
              ("SQL", 84, SUCCESS),
              ("System Design", 33, DANGER),
              ("Python", 71, PRIMARY)]
    y = Inches(3.5)
    for t, val, c in topics:
        add_text(s, Inches(6.9), y, Inches(2.2), Inches(0.32),
                 t, size=11, color=TEXT)
        # bar bg
        add_rect(s, Inches(8.9), y + Inches(0.08), Inches(3.4),
                 Inches(0.18), fill=BG_PANEL_2)
        # bar fill
        add_rect(s, Inches(8.9), y + Inches(0.08),
                 Inches(3.4 * val / 100.0), Inches(0.18), fill=c)
        add_text(s, Inches(12.35), y, Inches(0.5), Inches(0.32),
                 f"{val}", size=11, bold=True, color=TEXT,
                 align=PP_ALIGN.RIGHT)
        y += Inches(0.42)

    add_text(s, Inches(6.9), Inches(6.4), Inches(6.0), Inches(0.6),
             "Every quiz, diagnostic, and interview writes back into this "
             "vector — driving the next question's difficulty and the next "
             "interview's topic mix.",
             size=10.5, italic=True, color=MUTED, line_spacing=1.4)

    add_footer(s, 9, TOTAL, "Methodology · Profile")


# ---------------------------------------------------------------------------
# Slide 10 — Methodology: Company Intelligence
# ---------------------------------------------------------------------------
def slide_company_intel():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "09 · Methodology",
               "Company Intelligence — Perplexity → Gemini → cached.")

    # Pipeline 4-step horizontal
    steps = [
        ("User picks", "Company + Role.\n(Google × ML Engineer)", PRIMARY),
        ("Perplexity", "Sonar fetches live: Glassdoor, "
         "LeetCode, blogs, with citations.", ACCENT),
        ("Gemini synth", "Reshapes raw text into a ranked "
         "topic list + interview pattern summary.", SUCCESS),
        ("DB cache",  "CompanyInsights row keyed by "
         "(company, role) — never re-fetched.", WARNING),
    ]
    cw = Inches(3.0); ch = Inches(2.2)
    x0, y0 = Inches(0.7), Inches(2.4); gap = Inches(0.12)
    for i, (h, d, c) in enumerate(steps):
        x = x0 + i * (cw + gap)
        card(s, x, y0, cw, ch, fill=BG_PANEL, accent=c)
        # step number
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.25), y0 + Inches(0.2),
                                    Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = c
        circle.line.fill.background(); circle.shadow.inherit = False
        add_text(s, x + Inches(0.25), y0 + Inches(0.25),
                 Inches(0.55), Inches(0.5),
                 f"{i+1}", size=18, bold=True, color=TEXT,
                 align=PP_ALIGN.CENTER, font=HEAD_FONT)
        add_text(s, x + Inches(0.9), y0 + Inches(0.28),
                 cw - Inches(1.0), Inches(0.4), h,
                 size=13, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.25), y0 + Inches(0.95),
                 cw - Inches(0.4), Inches(1.2), d,
                 size=11, color=MUTED, line_spacing=1.4)
        # arrow between cards
        if i < len(steps) - 1:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                     x + cw - Inches(0.05),
                                     y0 + ch / 2 - Inches(0.15),
                                     Inches(0.22), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = DIM
            arr.line.fill.background(); arr.shadow.inherit = False

    # Output: ultra-personalized
    card(s, Inches(0.7), Inches(5.0), Inches(12.0), Inches(1.55),
         fill=BG_PANEL_2, accent=PRIMARY)
    add_text(s, Inches(0.9), Inches(5.15), Inches(11.5), Inches(0.4),
             "Output → Ultra-personalized plan",
             size=14, bold=True, color=ACCENT, font=HEAD_FONT)
    add_text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(1.0),
             "Intersection of  (a) user's diagnostic weaknesses  ∩  "
             "(b) company's high-frequency topics  →  prioritized 5–10 "
             "topic study queue + a targeted mock-interview that draws "
             "ONLY from that intersection. Library grows with every new "
             "company a user searches.",
             size=12, color=TEXT, line_spacing=1.45)

    add_footer(s, 10, TOTAL, "Company Intelligence")


# ---------------------------------------------------------------------------
# Slide 11 — Problems We Solved
# ---------------------------------------------------------------------------
def slide_problems_solved():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "10 · Engineering · Problems We Solved",
               "Hard problems shipped — not just talked about.")

    items = [
        ("Gemini JSON drift",
         "Multi-pass JSON parser (`_safe_parse_json` — 5 passes) handles "
         "trailing commas, markdown fences, partial outputs.", PRIMARY),
        ("Path personalization with no history",
         "Hard-coded Green seeds per role + Gemini-generated Yellow "
         "extensions; cold-start solved without a recommender.", ACCENT),
        ("Adaptive without infinite question banks",
         "Progressive easy → intermediate → advanced ladder, stop on "
         "first fail. Classifies a topic in ≤3 questions.", SUCCESS),
        ("Stale company data",
         "Perplexity (live) → Gemini (synth) → DB cache. Fresh on first "
         "hit, instant thereafter; library grows organically.", WARNING),
        ("Real-time skill updates",
         "Every quiz / interview event writes to UserSkillProfile with "
         "history[]; trend (improve/decline/stable) computed on read.", PRIMARY),
        ("Behavioral signals without multi-camera",
         "face-api.js 68-point landmarks → head-pose / gaze proxy for "
         "posture today; schema future-proofed for multi-camera.", ACCENT),
        ("Body + voice analysis cost",
         "Run AFTER the interview, async, from stored transcript + "
         "behavioral_stats — never blocks the live session.", SUCCESS),
        ("Follow-ups that actually feel adaptive",
         "Strategy classifier (vague/wrong/good/extend) routes to one "
         "of 4 Gemini prompts — answer-aware, not generic.", WARNING),
    ]
    cols = 2
    cw = Inches(6.05); ch = Inches(1.05)
    x0, y0 = Inches(0.7), Inches(2.3); gx, gy = Inches(0.15), Inches(0.12)
    for i, (h, d, c) in enumerate(items):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(s, x, y, cw, ch, fill=BG_PANEL, accent=c, accent_w=0.06)
        add_text(s, x + Inches(0.22), y + Inches(0.12),
                 cw - Inches(0.35), Inches(0.35), h,
                 size=12, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.22), y + Inches(0.45),
                 cw - Inches(0.35), Inches(0.6), d,
                 size=10.5, color=MUTED, line_spacing=1.35)

    add_footer(s, 11, TOTAL, "Problems Solved")


# ---------------------------------------------------------------------------
# Slide 12 — Tech Stack
# ---------------------------------------------------------------------------
def slide_tech_stack():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "11 · Tech Stack",
               "Best-in-class libraries chosen for adaptive UX.")

    columns = [
        ("Frontend", PRIMARY, [
            "React 18 + Vite",
            "shadcn/ui + @radix-ui/*",
            "Framer Motion",
            "@dnd-kit/core + sortable",
            "Recharts + SVG rings",
            "TanStack Query",
            "cmdk · vaul · resizable-panels",
            "react-markdown + remark-gfm",
            "lucide-react · react-hot-toast",
        ]),
        ("Backend", ACCENT, [
            "FastAPI (Python 3.11)",
            "Pydantic v2 schemas",
            "SQLAlchemy + Alembic",
            "PostgreSQL (prod) / SQLite (dev)",
            "JWT auth + role guards",
            "pdfplumber (resume OCR)",
            "httpx (async HTTP)",
            "Redis (rate limits, cache)",
            "Uvicorn / Gunicorn",
        ]),
        ("AI / Agentic", SUCCESS, [
            "Gemini 2.5 Flash (JSON mode)",
            "Perplexity Sonar (live web)",
            "OpenAI Whisper (ASR)",
            "face-api.js (gaze + 68 pts)",
            "Custom prompt registry",
            "5-pass JSON parser",
            "Strategy-routed follow-ups",
            "Async post-hoc analysis",
        ]),
        ("DevOps & Obs", WARNING, [
            "Docker · docker-compose",
            "GitHub Actions CI",
            "Vercel (frontend)",
            "Render / Fly (backend)",
            "Sentry (error tracking)",
            "PostHog (product analytics)",
            "Cloudflare R2 (assets)",
            ".env-driven config",
        ]),
    ]
    cw = Inches(3.05); ch = Inches(4.4)
    x0, y0 = Inches(0.7), Inches(2.3); gap = Inches(0.12)
    for i, (head, c, items) in enumerate(columns):
        x = x0 + i * (cw + gap)
        card(s, x, y0, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.25), y0 + Inches(0.2),
                 cw - Inches(0.4), Inches(0.4), head,
                 size=14, bold=True, color=c, font=HEAD_FONT)
        yy = y0 + Inches(0.65)
        for it in items:
            add_rect(s, x + Inches(0.27), yy + Inches(0.12),
                     Inches(0.08), Inches(0.08), fill=c)
            add_text(s, x + Inches(0.42), yy, cw - Inches(0.55),
                     Inches(0.38), it, size=10.5, color=TEXT,
                     line_spacing=1.25)
            yy += Inches(0.38)

    add_footer(s, 12, TOTAL, "Tech Stack")


# ---------------------------------------------------------------------------
# Slide 13 — Data Model
# ---------------------------------------------------------------------------
def slide_data_model():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "12 · Data Model",
               "Six tables — designed for adaptive learning + replay.")

    tables = [
        ("User",
         ["id, email, password_hash",
          "target_role, onboarding_complete",
          "created_at, last_active"], PRIMARY),
        ("UserSkillProfile",
         ["user_id, topic, score (0–100)",
          "confidence, last_updated",
          "history[] (trend signal)"], SUCCESS),
        ("LearningPath",
         ["user_id, role",
          "green_topics[], yellow_topics[]",
          "time_horizon, status"], ACCENT),
        ("InterviewSession",
         ["user_id, mode, topics_covered[]",
          "transcript JSON, report JSON",
          "behavioral_stats (extensible)"], WARNING),
        ("CompanyInsights",
         ["company × role (unique)",
          "topics[], patterns, summary",
          "source_data, analyzed_at"], RGBColor(0x06, 0xB6, 0xD4)),
        ("TopicArticle",
         ["(topic, job_role) key",
          "markdown_body (Gemini)",
          "generated_at"], DANGER),
    ]
    cols = 3
    cw = Inches(4.0); ch = Inches(2.0)
    x0, y0 = Inches(0.7), Inches(2.4); gx, gy = Inches(0.15), Inches(0.18)
    for i, (h, fields, c) in enumerate(tables):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(s, x, y, cw, ch, fill=BG_PANEL, accent=c)
        add_text(s, x + Inches(0.25), y + Inches(0.2),
                 cw - Inches(0.4), Inches(0.4),
                 h, size=14, bold=True, color=c, font=HEAD_FONT)
        yy = y + Inches(0.65)
        for f in fields:
            add_rect(s, x + Inches(0.27), yy + Inches(0.12),
                     Inches(0.07), Inches(0.07), fill=c)
            add_text(s, x + Inches(0.42), yy, cw - Inches(0.55),
                     Inches(0.4), f, size=10.5, color=TEXT,
                     font="Consolas", line_spacing=1.25)
            yy += Inches(0.36)

    add_footer(s, 13, TOTAL, "Data Model")


# ---------------------------------------------------------------------------
# Slide 14 — Body Language & Communication Analysis
# ---------------------------------------------------------------------------
def slide_body_comm():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "13 · Behavioral Analysis",
               "Body language + communication — async, after the interview.")

    # Left: signals
    add_text(s, Inches(0.7), Inches(2.4), Inches(6.0), Inches(0.4),
             "SIGNALS CAPTURED  (live, low-cost)",
             size=12, bold=True, color=ACCENT, font=HEAD_FONT)

    signals = [
        ("Gaze tracking",       "Off-camera % over time", PRIMARY),
        ("68-point landmarks",  "Head pose proxy for posture", PRIMARY),
        ("Expression timeline", "Calm / Stressed / Confident", ACCENT),
        ("Proctor flags",       "Multiple faces, looking-away", DANGER),
        ("Audio capture",       "Whisper transcript + timestamps", SUCCESS),
    ]
    y = Inches(2.85)
    for h, d, c in signals:
        card(s, Inches(0.7), y, Inches(6.0), Inches(0.6),
             fill=BG_PANEL, accent=c, accent_w=0.05)
        add_text(s, Inches(0.9), y + Inches(0.13),
                 Inches(2.4), Inches(0.32),
                 h, size=11.5, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, Inches(3.3), y + Inches(0.13),
                 Inches(3.4), Inches(0.32),
                 d, size=10.5, color=MUTED)
        y += Inches(0.66)

    # Right: post-hoc scores
    add_text(s, Inches(7.0), Inches(2.4), Inches(5.8), Inches(0.4),
             "POST-INTERVIEW REPORT  (async)",
             size=12, bold=True, color=ACCENT, font=HEAD_FONT)
    posthoc = [
        ("Speaking pace",     "WPM from transcript ÷ duration", PRIMARY),
        ("Filler words",      "um / uh / like / you-know count", PRIMARY),
        ("Vocabulary richness","Gemini lexical-diversity score", ACCENT),
        ("Grammar & fluency", "Gemini analysis with quoted moments", ACCENT),
        ("Eye contact %",     "From gaze timeline", SUCCESS),
        ("Posture proxy",     "Head-tilt + face-position heuristics", SUCCESS),
    ]
    y = Inches(2.85)
    for h, d, c in posthoc:
        card(s, Inches(7.0), y, Inches(5.8), Inches(0.55),
             fill=BG_PANEL_2, accent=c, accent_w=0.05)
        add_text(s, Inches(7.2), y + Inches(0.1),
                 Inches(2.4), Inches(0.32),
                 h, size=11.5, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, Inches(9.6), y + Inches(0.1),
                 Inches(3.1), Inches(0.32),
                 d, size=10.5, color=MUTED)
        y += Inches(0.62)

    add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
             "Architecture is extensible — behavioral_stats is a JSON column "
             "ready for multi-camera body-keypoint streams.",
             size=10.5, italic=True, color=MUTED)

    add_footer(s, 14, TOTAL, "Behavioral Analysis")


# ---------------------------------------------------------------------------
# Slide 15 — Roadmap & Phases
# ---------------------------------------------------------------------------
def slide_roadmap():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    section_band(s)
    add_header(s, "14 · Roadmap",
               "Phased delivery — value at every step.")

    phases = [
        ("P0", "Onboarding · Resume OCR · Path Configurator",
         "Role pick, PDF→Gemini analysis, Green/Yellow drag-drop.", PRIMARY),
        ("P1", "Adaptive Diagnostic · Articles · Mini-Quizzes",
         "Progressive Q ladder, cached articles, 5-Q quizzes, "
         "skill-meter writes.", ACCENT),
        ("P2", "Time-Based + Company-Specific Personalization",
         "24h / 1w / 1m / 3m / 6m re-ranking. Perplexity company "
         "intelligence with cache.", SUCCESS),
        ("P3", "Body Language + Communication Analytics",
         "Gaze, expression, voice pace, filler words, fluency — "
         "async post-interview report.", WARNING),
        ("P4", "Full-Syllabus Interviews + Trend Reports",
         "Quick full interview mode + history compare across "
         "past sessions.", RGBColor(0x06, 0xB6, 0xD4)),
    ]
    # vertical timeline
    x_line = Inches(1.4)
    add_rect(s, x_line, Inches(2.35), Inches(0.03), Inches(4.4),
             fill=DIM)
    y = Inches(2.35)
    for label, h, d, c in phases:
        # dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x_line - Inches(0.18), y + Inches(0.05),
                                 Inches(0.4), Inches(0.4))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background(); dot.shadow.inherit = False
        add_text(s, x_line - Inches(0.18), y + Inches(0.08),
                 Inches(0.4), Inches(0.34), label,
                 size=10, bold=True, color=TEXT,
                 align=PP_ALIGN.CENTER, font=HEAD_FONT)
        # card
        card(s, Inches(2.0), y, Inches(10.7), Inches(0.75),
             fill=BG_PANEL, accent=c, accent_w=0.05)
        add_text(s, Inches(2.2), y + Inches(0.1),
                 Inches(4.0), Inches(0.35),
                 h, size=13, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, Inches(6.2), y + Inches(0.1),
                 Inches(6.4), Inches(0.55),
                 d, size=10.5, color=MUTED, line_spacing=1.3)
        y += Inches(0.9)

    add_footer(s, 15, TOTAL, "Roadmap")


# ---------------------------------------------------------------------------
# Slide 16 — Differentiators / Close
# ---------------------------------------------------------------------------
def slide_close():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG_DEEP)
    section_band(s)

    add_text(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.4),
             "15 · WHY THIS WINS",
             size=12, bold=True, color=ACCENT, font=HEAD_FONT)
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(1.0),
             "An interview coach that learns you faster than you "
             "learn the syllabus.",
             size=30, bold=True, color=TEXT, font=HEAD_FONT,
             line_spacing=1.1)

    diff = [
        ("Per-user skill vector",
         "Compounds across every activity — no cold-start after day 1."),
        ("Two-AI orchestration",
         "Gemini for reasoning, Perplexity for freshness — playing to strengths."),
        ("Socratic, not scripted",
         "Follow-ups respond to *this* answer, not a template."),
        ("Holistic interview signal",
         "Tech + behavioral + body + voice in one report."),
        ("Time- & company-aware",
         "Replans the syllabus for a 24h crunch or a Google ML-E role."),
        ("Extensible architecture",
         "Multi-camera, more models, more roles — slot in without schema breaks."),
    ]
    cols = 3
    cw = Inches(4.05); ch = Inches(1.45)
    x0, y0 = Inches(0.7), Inches(3.0); gx, gy = Inches(0.12), Inches(0.18)
    for i, (h, d) in enumerate(diff):
        row, col = divmod(i, cols)
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(s, x, y, cw, ch, fill=BG_PANEL,
             accent=[PRIMARY, ACCENT, SUCCESS, WARNING,
                     RGBColor(0x06, 0xB6, 0xD4), DANGER][i],
             accent_w=0.06)
        add_text(s, x + Inches(0.22), y + Inches(0.18),
                 cw - Inches(0.4), Inches(0.4),
                 h, size=13, bold=True, color=TEXT, font=HEAD_FONT)
        add_text(s, x + Inches(0.22), y + Inches(0.55),
                 cw - Inches(0.4), Inches(0.85),
                 d, size=10.5, color=MUTED, line_spacing=1.35)

    add_text(s, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
             "Thank you  ·  Questions?",
             size=14, bold=True, color=ACCENT, font=HEAD_FONT)
    add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
             "Nitin  ·  InterviewVault  ·  2026",
             size=10, color=DIM, italic=True)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
slide_title()
slide_problem_statements()
slide_vision()
slide_system_arch()
slide_backend()
slide_frontend()
slide_ai_layer()
slide_methodology_path()
slide_methodology_socratic()
slide_company_intel()
slide_problems_solved()
slide_tech_stack()
slide_data_model()
slide_body_comm()
slide_roadmap()
slide_close()

out = "InterviewVault_Master_Deck.pptx"
prs.save(out)
print(f"[OK] Wrote {out}  —  {sum(1 for _ in prs.slides)} slides")
